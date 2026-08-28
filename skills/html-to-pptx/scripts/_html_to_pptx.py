"""HTML slide / poster -> editable PPTX (semantic reconstruction).

Design goal: **every visual element becomes its own PPTX object.** A deck that
opens in PowerPoint should let you grab the arrow, the dragon and the icon
independently — not find them fused into one flat backdrop picture.

Pipeline (deterministic; a single optional LLM call may only re-label roles):

  1. extract    Playwright renders the canvas (``.slide-canvas`` / ``.poster-canvas``
                / ``[data-canvas-width]``, auto-detected) and walks the DOM,
                collecting each element's canvas-relative box, computed style,
                rotation, stacking order and whether it owns renderable content.
                Coordinates come from ``getBoundingClientRect`` so grid/flex
                layouts are captured precisely.
  2. classify   A deterministic four-way rule assigns a role:
                  text     editable text box (plus its own fill, if painted)
                  shape    native (rounded) rectangle; children render on top
                  picture  rasterized **on its own transparent canvas** — images,
                           inline SVG, and painted leaves too complex for a
                           native shape (partial borders, gradients, rotation)
                  image    not reconstructed; stays in the residual backdrop
  3. rasterize  Each ``picture`` is screenshotted in isolation (every other
                element hidden, transparent background) so it lands in the deck
                as a standalone, movable picture shape.
  4. backdrop   Everything reconstructed is neutralized and the canvas is
                re-shot. Only what we failed to rebuild survives, so the deck is
                never missing pixels — usually just the page background.
  5. build      python-pptx renders backdrop first, then every element in
                stacking order (z-index, then document order).
  6. verify     The saved file is re-opened to confirm shape/text/picture counts.

Slide dimensions follow the canvas aspect ratio (long side 13.333in), so a
1920x1080 slide stays 16:9 and a 1067x1600 poster stays portrait. The scale is
uniform on both axes, so nothing is ever stretched.
"""

from __future__ import annotations

import json
import math
import os
import re
import threading
from contextlib import contextmanager
from functools import partial
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from typing import Any, Dict, Iterator, List, Optional, Tuple
from urllib.parse import quote

# Long edge of the generated slide. 1920x1080 -> 13.333in x 7.5in (the classic
# 16:9 deck); a portrait poster keeps its own ratio against the same long edge.
LONG_SIDE_IN = 13.3333333
EMU_PER_INCH = 914400
EMU_PER_PT = 12700
BOLD_WEIGHT_THRESHOLD = 600
# An element covering most of the canvas is treated as backdrop, never a shape.
MAX_SHAPE_AREA_FRACTION = 0.9
# Minimum opaque alpha for a color to count as "painted".
MIN_VISIBLE_ALPHA = 0.02
# Below this angle a rotation is noise from sub-pixel layout, not design intent.
MIN_ROTATION_DEG = 0.5

# Kept for backwards compatibility with callers that imported the 16:9 defaults.
DEFAULT_SLIDE_WIDTH_EMU = 12192000
DEFAULT_SLIDE_HEIGHT_EMU = 6858000

CANVAS_SELECTORS = (".slide-canvas", ".poster-canvas", "[data-canvas-width]", "#poster")

_CJK_RE = re.compile(r"[\u3400-\u9fff\uf900-\ufaff\uff00-\uffef]")
_RGBA_RE = re.compile(
    r"rgba?\(\s*([0-9.]+)[,\s]+([0-9.]+)[,\s]+([0-9.]+)(?:[,\s/]+([0-9.]+%?))?\s*\)"
)

# Collected once in the page. ``data-wca-id`` is stamped on every element so the
# rasterize and neutralize passes can re-select the exact same nodes.
_EXTRACT_JS = r"""
(canvasSel) => {
  const canvas = document.querySelector(canvasSel);
  if (!canvas) return null;
  const cb = canvas.getBoundingClientRect();
  const SKIP = new Set(['script','style','defs','clippath','marker','filter','lineargradient','radialgradient','mask','symbol','title','metadata']);
  const hasDirectText = (node) => {
    for (const n of node.childNodes) if (n.nodeType === 3 && n.textContent.trim()) return true;
    return false;
  };
  const parseAlpha = (c) => {
    const m = c && c.match(/rgba?\(\s*[\d.]+[,\s]+[\d.]+[,\s]+[\d.]+(?:[,\s\/]+([\d.]+%?))?/);
    if (!m || m[1] == null) return c && c !== 'transparent' ? 1 : 0;
    return m[1].endsWith('%') ? parseFloat(m[1]) / 100 : parseFloat(m[1]);
  };
  const sides = ['Top','Right','Bottom','Left'];
  // "Painted" = contributes ink of its own (fill, border, gradient or shadow).
  const painted = (st) => {
    if (parseAlpha(st.backgroundColor) > 0.02) return true;
    if ((st.backgroundImage || 'none') !== 'none') return true;
    if ((st.boxShadow || 'none') !== 'none') return true;
    for (const s of sides) {
      if (parseFloat(st['border' + s + 'Width']) > 0
          && st['border' + s + 'Style'] !== 'none'
          && parseAlpha(st['border' + s + 'Color']) > 0.02) return true;
    }
    return false;
  };
  const insideSvg = (el) => { const a = el.closest('svg'); return !!a && a !== el; };
  const visible = (el) => {
    const st = getComputedStyle(el);
    if (st.display === 'none' || st.visibility === 'hidden') return false;
    if (parseFloat(st.opacity || '1') <= 0.01) return false;
    return true;
  };
  // Does this node put anything on screen by itself?
  const renderable = (el) => {
    const t = el.tagName.toLowerCase();
    if (t === 'img' || t === 'svg' || t === 'canvas' || t === 'video') return true;
    if (hasDirectText(el)) return true;
    return painted(getComputedStyle(el));
  };
  const rotationOf = (st) => {
    const tr = st.transform;
    if (!tr || tr === 'none') return 0;
    const m = tr.match(/matrix\(([^)]+)\)/);
    if (!m) return 0;
    const p = m[1].split(',').map(Number);
    return Math.atan2(p[1], p[0]) * 180 / Math.PI;
  };
  // Rotation is inherited through ancestor transforms exactly as scale is. A pill
  // tilted by a wrapper reports `transform: none` on itself, so reading only its
  // own style exports it bolt upright.
  const contentRotationOf = (el) => {
    let deg = 0;
    for (let n = el; n && n !== canvas.parentElement; n = n.parentElement) {
      deg += rotationOf(getComputedStyle(n));
    }
    return deg;
  };
  // Cumulative scale from every transform between the element and the canvas,
  // its own included. Boxes come from getBoundingClientRect (already scaled) but
  // font size and border width come from computed style (not scaled), so without
  // this factor text inside a scaled wrapper renders 1/scale too large.
  const contentScaleOf = (el) => {
    let s = 1;
    for (let n = el; n && n !== canvas.parentElement; n = n.parentElement) {
      const t = getComputedStyle(n).transform;
      if (!t || t === 'none') continue;
      const m = t.match(/matrix\(([^)]+)\)/);
      if (!m) continue;
      const p = m[1].split(',').map(Number);
      s *= (Math.hypot(p[0], p[1]) + Math.hypot(p[2], p[3])) / 2;
    }
    return s > 0 ? s : 1;
  };
  // Cumulative CSS opacity, own included. Elements are extracted individually, so
  // a faded ancestor has to be folded into each descendant's own alpha.
  const contentOpacityOf = (el) => {
    let o = 1;
    for (let n = el; n && n !== canvas.parentElement; n = n.parentElement) {
      const v = parseFloat(getComputedStyle(n).opacity || '1');
      if (!isNaN(v)) o *= v;
    }
    return o;
  };
  // Visible slice after every clipping ancestor. An element sticking out of an
  // overflow:hidden crop must not be drawn full-size in the deck.
  const clipRectOf = (el, r) => {
    let x0 = r.left, y0 = r.top, x1 = r.right, y1 = r.bottom;
    for (let n = el.parentElement; n && n !== canvas.parentElement; n = n.parentElement) {
      const st = getComputedStyle(n);
      const hides = [st.overflow, st.overflowX, st.overflowY].some(
        (v) => v === 'hidden' || v === 'clip' || v === 'auto' || v === 'scroll');
      if (!hides) continue;
      const nr = n.getBoundingClientRect();
      x0 = Math.max(x0, nr.left); y0 = Math.max(y0, nr.top);
      x1 = Math.min(x1, nr.right); y1 = Math.min(y1, nr.bottom);
    }
    return { left: x0, top: y0, width: Math.max(0, x1 - x0), height: Math.max(0, y1 - y0) };
  };
  // z-index only orders siblings inside one stacking context, and it is ignored
  // on a static element that is not a flex/grid item. Reading it as a global key
  // lets a card with `z-index:1` sort above the text it contains.
  const stackZ = (el) => {
    const st = getComputedStyle(el);
    const z = parseInt(st.zIndex, 10);
    if (isNaN(z)) return 0;
    if (st.position !== 'static') return z;
    const par = el.parentElement;
    if (par && /(flex|grid)$/.test(getComputedStyle(par).display)) return z;
    return 0;
  };
  const sibIdx = new WeakMap();
  const siblingIndex = (el) => {
    if (sibIdx.has(el)) return sibIdx.get(el);
    const par = el.parentElement;
    if (!par) return 0;
    let i = 0;
    for (const kid of par.children) sibIdx.set(kid, i++);
    return sibIdx.get(el) ?? 0;
  };
  // Paint order as a path from the canvas down: each step contributes the node's
  // effective z-index then its position among siblings. Comparing two paths
  // lexicographically resolves them at their lowest common ancestor, and an
  // ancestor's path is a prefix of its descendant's, so a parent always paints
  // first — which is exactly the rule a flat z-index sort was breaking.
  const paintPathOf = (el) => {
    const chain = [];
    for (let n = el; n && n !== canvas; n = n.parentElement) chain.push(n);
    const out = [];
    for (let i = chain.length - 1; i >= 0; i--) {
      out.push(stackZ(chain[i]), siblingIndex(chain[i]));
    }
    return out;
  };
  const hasPseudo = (el) => {
    for (const p of ['::before', '::after']) {
      const cs = getComputedStyle(el, p);
      if (!cs) continue;
      const c = cs.content;
      if (c && c !== 'none' && c !== 'normal' && parseFloat(cs.width || '0') >= 0 && painted(cs)) return true;
      if (c && c !== 'none' && c !== 'normal' && c !== '""') return true;
    }
    return false;
  };
  const resolveColor = (node) => {
    const s = getComputedStyle(node);
    let c = (s.webkitTextFillColor && s.webkitTextFillColor !== 'currentcolor')
      ? s.webkitTextFillColor : s.color;
    if (parseAlpha(c) < 0.1) {
      const g = (s.backgroundImage || '').match(/rgba?\([^)]*\)/);
      if (g) return g[0];
      if (parseAlpha(s.color) >= 0.1) return s.color;
    }
    return c;
  };
  const isInline = (node) => {
    const d = getComputedStyle(node).display;
    return d.startsWith('inline') || d === 'contents';
  };
  // OOXML names exactly one typeface, so the CSS chain has to collapse to the
  // family the browser actually painted with. Generic keywords are skipped: they
  // resolve per-platform and carry no name PowerPoint could look up.
  const GENERIC = /^(sans-serif|serif|monospace|cursive|fantasy|system-ui|ui-sans-serif|ui-serif|ui-monospace|ui-rounded|math|emoji|fangsong|-apple-system|blinkmacsystemfont)$/i;
  // document.fonts.check() answers "can this string be rendered", which is always
  // true once fallback kicks in, so it cannot tell an installed family from a
  // missing one. Measuring glyph widths against baseline families can.
  const PROBE = 'mmmmmmmmmmlliWWWW@@@@国永';
  const measureCtx = document.createElement('canvas').getContext('2d');
  const baseWidths = {};
  for (const b of ['monospace', 'serif', 'sans-serif']) {
    measureCtx.font = '72px ' + b;
    baseWidths[b] = measureCtx.measureText(PROBE).width;
  }
  const installed = (name) => {
    for (const b of ['monospace', 'serif', 'sans-serif']) {
      measureCtx.font = '72px "' + name + '", ' + b;
      if (measureCtx.measureText(PROBE).width !== baseWidths[b]) return true;
    }
    return false;
  };
  const familyCache = new Map();
  const resolveFamily = (chain) => {
    if (familyCache.has(chain)) return familyCache.get(chain);
    let picked = '';
    for (const raw of (chain || '').split(',')) {
      const name = raw.trim().replace(/^["']|["']$/g, '');
      if (!name || GENERIC.test(name)) continue;
      if (!picked) picked = name;   // keep the first concrete name as a fallback
      let ok = false;
      try { ok = installed(name); } catch (e) { /* malformed family name */ }
      if (ok) { picked = name; break; }   // this is the one actually painted
    }
    familyCache.set(chain, picked);
    return picked;
  };
  const spacingOf = (cs) => {
    const v = parseFloat(cs.letterSpacing);
    return isNaN(v) ? 0 : v;               // 'normal' -> 0
  };
  // Flatten a text element into lines of styled runs so multi-color titles and
  // muted suffixes keep their per-span color / weight / size.
  const buildRich = (root) => {
    const lines = [[]];
    const cur = () => lines[lines.length - 1];
    const walk = (node) => {
      for (const ch of node.childNodes) {
        if (ch.nodeType === 3) {
          const t = ch.textContent.replace(/\s+/g, ' ');
          if (t) {
            const cs = getComputedStyle(node);
            cur().push({ text: t, color: resolveColor(node), weight: cs.fontWeight || '400',
              size: parseFloat(cs.fontSize) || 0, family: resolveFamily(cs.fontFamily),
              italic: cs.fontStyle === 'italic' || cs.fontStyle === 'oblique',
              spacing: spacingOf(cs) });
          }
        } else if (ch.nodeType === 1) {
          const t = ch.tagName.toLowerCase();
          if (SKIP.has(t)) continue;
          if (t === 'br') { lines.push([]); continue; }
          const cs = getComputedStyle(ch);
          if (cs.display === 'none' || cs.visibility === 'hidden') continue;
          const block = !isInline(ch);
          if (block && cur().length) lines.push([]);
          walk(ch);
          if (block && cur().length) lines.push([]);
        }
      }
    };
    walk(root);
    const cleaned = [];
    for (const line of lines) {
      const runs = line.filter((r) => r.text !== '');
      if (runs.length && runs.some((r) => r.text.trim() !== '')) {
        runs[0].text = runs[0].text.replace(/^\s+/, '');
        runs[runs.length - 1].text = runs[runs.length - 1].text.replace(/\s+$/, '');
        cleaned.push(runs.filter((r) => r.text !== ''));
      }
    }
    return cleaned;
  };

  // A run of text carries no geometry of its own, so a row whose chunks are
  // pushed apart by flex `gap` or by margins on inline separators collapses to a
  // fraction of its width once merged into one paragraph. Wrapping each bare
  // text run in a span makes every chunk a measurable element that lands at its
  // own position instead. Layout is unaffected: an inline span around a text run
  // occupies exactly the box the anonymous run already did.
  const spacedRow = (box) => {
    const disp = getComputedStyle(box).display;
    if (/^(inline-)?(flex|grid)$/.test(disp)) {
      const gap = parseFloat(getComputedStyle(box).gap);
      if (!isNaN(gap) && gap > 0.5) return true;
    }
    for (const kid of box.children) {
      const ks = getComputedStyle(kid);
      for (const prop of ['marginLeft', 'marginRight', 'paddingLeft', 'paddingRight']) {
        if (Math.abs(parseFloat(ks[prop]) || 0) > 0.5) return true;
      }
    }
    return false;
  };
  for (const box of canvas.querySelectorAll('*')) {
    if (!box.children.length || !hasDirectText(box) || !spacedRow(box)) continue;
    for (const node of [...box.childNodes]) {
      if (node.nodeType !== 3 || !node.textContent.trim()) continue;
      const span = document.createElement('span');
      span.setAttribute('data-wca-anon', '1');
      node.parentNode.replaceChild(span, node);
      span.appendChild(node);
    }
  }

  const out = [];
  let idx = 0;
  for (const el of canvas.querySelectorAll('*')) {
    const tag = el.tagName.toLowerCase();
    if (SKIP.has(tag)) continue;
    if (insideSvg(el)) continue;          // svg internals ride along with their root
    if (!visible(el)) continue;
    const r = el.getBoundingClientRect();
    if (r.width <= 0.5 || r.height <= 0.5) continue;
    const vis = clipRectOf(el, r);
    if (vis.width <= 0.5 || vis.height <= 0.5) continue;   // clipped away entirely

    const st = getComputedStyle(el);
    let textAncestor = false;
    for (let a = el.parentElement; a && a !== canvas; a = a.parentElement) {
      if (hasDirectText(a)) { textAncestor = true; break; }
    }
    let renderableDescendant = false;
    for (const d of el.querySelectorAll('*')) {
      if (SKIP.has(d.tagName.toLowerCase())) continue;
      if (insideSvg(d) || !visible(d)) continue;
      const dr = d.getBoundingClientRect();
      if (dr.width <= 0.5 || dr.height <= 0.5) continue;
      if (renderable(d)) { renderableDescendant = true; break; }
    }
    const bw = sides.map((s) => parseFloat(st['border' + s + 'Width']) || 0);
    const bc = sides.map((s) => st['border' + s + 'Color']);
    const uniformBorder = bw.every((v) => Math.abs(v - bw[0]) < 0.01)
      && bc.every((v) => v === bc[0]);
    const radii = ['TopLeft','TopRight','BottomLeft','BottomRight']
      .map((c) => parseFloat(st['border' + c + 'Radius']) || 0);
    const zi = st.zIndex;

    el.setAttribute('data-wca-id', String(idx));
    const direct = hasDirectText(el);
    out.push({
      id: idx,
      tag,
      className: (el.getAttribute('class') || '').slice(0, 140),
      x: r.left - cb.left, y: r.top - cb.top, w: r.width, h: r.height,
      vx: r.left, vy: r.top,                        // viewport coords, for screenshot clips
      visX: vis.left - cb.left, visY: vis.top - cb.top,
      visW: vis.width, visH: vis.height,
      visVX: vis.left, visVY: vis.top,
      clipped: (vis.width < r.width - 0.5) || (vis.height < r.height - 0.5),
      contentScale: contentScaleOf(el),
      contentOpacity: contentOpacityOf(el),
      layoutW: el.offsetWidth || r.width,           // pre-transform size (rotation aside)
      layoutH: el.offsetHeight || r.height,
      areaFrac: (r.width * r.height) / Math.max(1, cb.width * cb.height),
      isSvg: tag === 'svg',
      isImg: tag === 'img' || tag === 'canvas' || tag === 'video',
      hasDirectText: direct,
      hasTextAncestor: textAncestor,
      hasRenderableDescendant: renderableDescendant,
      hasPseudo: hasPseudo(el),
      rotation: contentRotationOf(el),
      zIndex: (zi === 'auto' || zi === '' || zi == null) ? 0 : (parseInt(zi, 10) || 0),
      paintPath: paintPathOf(el),
      innerText: direct ? el.innerText : '',
      richText: direct ? buildRich(el) : [],
      fontSize: parseFloat(st.fontSize) || 0,
      fontWeight: st.fontWeight || '400',
      color: st.color,
      textAlign: st.textAlign || 'start',
      mixBlendMode: st.mixBlendMode || 'normal',
      display: st.display || 'block',
      flexDirection: st.flexDirection || 'row',
      justifyContent: st.justifyContent || 'normal',
      justifyItems: st.justifyItems || 'normal',
      alignItems: st.alignItems || 'normal',
      writingMode: st.writingMode || 'horizontal-tb',
      fontFamily: resolveFamily(st.fontFamily),
      fontStyle: st.fontStyle || 'normal',
      letterSpacing: spacingOf(st),
      lineHeight: parseFloat(st.lineHeight) || 0,   // 'normal' -> 0, meaning auto
      opacity: parseFloat(st.opacity || '1'),
      bgColor: st.backgroundColor,
      bgImage: st.backgroundImage || 'none',
      borderRadius: Math.max.apply(null, radii),
      uniformRadius: radii.every((v) => Math.abs(v - radii[0]) < 0.01),
      borderWidth: bw[0],
      borderColor: bc[0],
      borderWidths: bw,                             // Top, Right, Bottom, Left
      borderColors: bc,
      uniformBorder: uniformBorder,
      boxShadow: st.boxShadow || 'none',
    });
    idx++;
  }
  return { canvas: { width: cb.width, height: cb.height }, elements: out };
}
"""

# Re-measure text under a neutral fallback font (Arial + PingFang/YaHei, present
# on typical macOS/Windows PowerPoint installs) to learn how much each text box
# must shrink to fit its ORIGINAL box when the design fonts are unavailable. The
# shrink is baked into the font size so it never depends on the opening app
# recomputing "shrink text on overflow" (PowerPoint does not, on open).
# Safety net for font substitution: re-measure each text box in the family that
# will actually be written to the pptx. When that family is the one the page
# already painted with, the ratio is 1 and nothing shrinks; it only bites when we
# had to fall back, where the substitute is wider and would overlap neighbours.
_MEASURE_FIT_JS = r"""
(items) => {
  const scales = {};
  const patched = [];
  for (const it of items) {
    const el = document.querySelector(`[data-wca-id="${it.id}"]`);
    if (!el) { scales[it.id] = 1; continue; }
    if (it.family) {
      patched.push([el, el.style.getPropertyValue('font-family'),
                    el.style.getPropertyPriority('font-family')]);
      el.style.setProperty('font-family', '"' + it.family + '"', 'important');
    }
  }
  void document.body.offsetHeight;  // force reflow
  for (const it of items) {
    const el = document.querySelector(`[data-wca-id="${it.id}"]`);
    if (!el) continue;
    const r = el.getBoundingClientRect();
    const wr = it.w > 0 && r.width > 0 ? it.w / r.width : 1;
    const hr = it.h > 0 && r.height > 0 ? it.h / r.height : 1;
    scales[it.id] = Math.max(0.3, Math.min(1, wr, hr));
  }
  for (const [el, value, prio] of patched) {
    if (value) el.style.setProperty('font-family', value, prio);
    else el.style.removeProperty('font-family');
  }
  return scales;
}
"""

# Isolation for per-element rasterization: strip page/canvas paint, hide every
# element, then re-show the target *and all of its descendants* (a parent set to
# `visible` cannot un-hide children that carry their own !important hidden).
_ISOLATE_JS = r"""
(payload) => {
  const canvas = document.querySelector(payload.canvasSel);
  if (!canvas) return;
  for (const node of [document.documentElement, document.body, canvas]) {
    node.style.setProperty('background', 'transparent', 'important');
    node.style.setProperty('background-image', 'none', 'important');
  }
  for (const el of canvas.querySelectorAll('*')) {
    el.style.setProperty('visibility', 'hidden', 'important');
  }
  const target = canvas.querySelector(`[data-wca-id="${payload.id}"]`);
  if (!target) return;
  target.style.setProperty('visibility', 'visible', 'important');
  for (const d of target.querySelectorAll('*')) {
    d.style.setProperty('visibility', 'visible', 'important');
  }
}
"""

_RESET_ISOLATION_JS = r"""
(canvasSel) => {
  const canvas = document.querySelector(canvasSel);
  if (!canvas) return;
  for (const node of [document.documentElement, document.body, canvas]) {
    node.style.removeProperty('background');
    node.style.removeProperty('background-image');
  }
  for (const el of canvas.querySelectorAll('*')) el.style.removeProperty('visibility');
}
"""

# Applied before the residual-backdrop shot. Each role gives up exactly the ink
# it contributes to the rebuilt deck, so whatever is left over is what we failed
# to reconstruct — and it still ships, inside the backdrop.
_NEUTRALIZE_JS = r"""
(payload) => {
  const strip = (el) => {
    el.style.setProperty('background', 'transparent', 'important');
    el.style.setProperty('background-image', 'none', 'important');
    el.style.setProperty('border', 'none', 'important');
    el.style.setProperty('box-shadow', 'none', 'important');
    el.style.setProperty('backdrop-filter', 'none', 'important');
    el.style.setProperty('-webkit-backdrop-filter', 'none', 'important');
  };
  for (const id of payload.shapeIds) {
    const el = document.querySelector(`[data-wca-id="${id}"]`);
    if (el) strip(el);                       // keep children: they render on top
  }
  for (const id of payload.textIds) {
    const el = document.querySelector(`[data-wca-id="${id}"]`);
    if (!el) continue;
    strip(el);
    el.style.setProperty('color', 'transparent', 'important');
    el.style.setProperty('-webkit-text-fill-color', 'transparent', 'important');
  }
  for (const id of payload.pictureIds) {
    const el = document.querySelector(`[data-wca-id="${id}"]`);
    if (!el) continue;                       // rasterized whole: hide the subtree
    el.style.setProperty('visibility', 'hidden', 'important');
    for (const d of el.querySelectorAll('*')) {
      d.style.setProperty('visibility', 'hidden', 'important');
    }
  }
}
"""


class _QuietHandler(SimpleHTTPRequestHandler):
    def log_message(self, _format: str, *args: Any) -> None:  # noqa: A003
        del args


@contextmanager
def _serve_preview(index: Path) -> Iterator[str]:
    """Serve ``index``'s directory on a random loopback port.

    Self-contained (mirrors quality.screenshots._serve_preview) so this export
    module does not boot the whole quality-gate/providers stack just to render.
    """
    root = index.parent
    handler = partial(_QuietHandler, directory=str(root))
    server = ThreadingHTTPServer(("127.0.0.1", 0), handler)
    server.daemon_threads = True
    thread = threading.Thread(
        target=server.serve_forever, name="pptx-preview-http", daemon=True
    )
    thread.start()
    relative_index = index.relative_to(root).as_posix()
    host, port = server.server_address[:2]
    url = f"http://{host}:{port}/{quote(relative_index)}"
    try:
        yield url
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=3)


def _parse_rgba(value: Optional[str]) -> Optional[Tuple[int, int, int, float]]:
    if not value:
        return None
    match = _RGBA_RE.search(value)
    if not match:
        return None
    r, g, b, a = match.groups()
    alpha = 1.0
    if a is not None:
        alpha = float(a[:-1]) / 100.0 if a.endswith("%") else float(a)
    return int(float(r)), int(float(g)), int(float(b)), alpha


def _alpha(value: Optional[str]) -> float:
    parsed = _parse_rgba(value)
    return parsed[3] if parsed else 0.0


def _gradient_first_color(bg_image: Optional[str]) -> Optional[Tuple[int, int, int]]:
    """First color stop of a CSS gradient (computed style is normalized to rgb()).

    Used to recover the intended color of gradient-clipped ("transparent") text,
    which a single-pixel sample would miss when it lands between glyphs.
    """
    if not bg_image or "gradient" not in bg_image:
        return None
    parsed = _parse_rgba(bg_image)
    return parsed[:3] if parsed else None


def _has_gradient(el: Dict[str, Any]) -> bool:
    return "gradient" in (el.get("bgImage") or "")


# CSS keyword directions, as the angle convention used by `linear-gradient`:
# 0deg points up and angles grow clockwise.
_GRADIENT_KEYWORD_DEG = {
    "to top": 0.0,
    "to right": 90.0,
    "to bottom": 180.0,
    "to left": 270.0,
    "to top right": 45.0,
    "to right top": 45.0,
    "to bottom right": 135.0,
    "to right bottom": 135.0,
    "to bottom left": 225.0,
    "to left bottom": 225.0,
    "to top left": 315.0,
    "to left top": 315.0,
}


def _split_top_level(text: str) -> List[str]:
    """Split on commas that are not inside parentheses (rgb(...) holds its own)."""
    parts: List[str] = []
    depth = 0
    start = 0
    for i, ch in enumerate(text):
        if ch == "(":
            depth += 1
        elif ch == ")":
            depth -= 1
        elif ch == "," and depth == 0:
            parts.append(text[start:i])
            start = i + 1
    parts.append(text[start:])
    return [p.strip() for p in parts if p.strip()]


def _linear_gradient(
    bg_image: Optional[str],
) -> Optional[Tuple[float, List[Tuple[int, int, int, float]]]]:
    """``(css_angle_deg, stops)`` for a lone CSS ``linear-gradient``, else None.

    Returns None for radial/conic gradients and for stacked multi-layer
    backgrounds — shapes PowerPoint cannot express, whose ink is better left in
    the backdrop than approximated by one flat color.
    """
    value = (bg_image or "").strip()
    if not value.startswith("linear-gradient(") or not value.endswith(")"):
        return None
    inner = value[len("linear-gradient(") : -1]
    tokens = _split_top_level(inner)
    if len(tokens) < 2:
        return None

    angle = 180.0  # CSS default is `to bottom`
    head = tokens[0].lower()
    if not _RGBA_RE.search(tokens[0]) and not tokens[0].startswith("#"):
        if head.endswith("deg"):
            try:
                angle = float(head[:-3])
            except ValueError:
                return None
        elif head in _GRADIENT_KEYWORD_DEG:
            angle = _GRADIENT_KEYWORD_DEG[head]
        else:
            return None  # turn/rad units, interpolation hints: not worth guessing
        tokens = tokens[1:]

    stops: List[Tuple[int, int, int, float]] = []
    for token in tokens:
        parsed = _parse_rgba(token)
        if parsed is None:
            return None
        stops.append(parsed)
    if len(stops) < 2:
        return None
    return angle, stops


_FLEX_DISPLAY_RE = re.compile(r"^(inline-)?(flex|grid)$")
# CSS box-alignment keywords, as the `text-align` value that looks the same.
_ALIGN_KEYWORDS = {
    "center": "center",
    "flex-end": "right",
    "end": "right",
    "right": "right",
    "flex-start": "left",
    "start": "left",
    "left": "left",
}


def _effective_align(el: Dict[str, Any]) -> str:
    """Horizontal alignment as rendered, not merely as ``text-align`` says.

    A pill that centers its label with ``justify-content`` leaves ``text-align``
    at the inherited ``start``, so reading that alone pins the words to the left
    edge of a box the browser draws them centered in.

    Distribution keywords (``space-between`` and friends) and ``normal`` fall
    through to ``text-align``: with a single run of text they put it at the
    start anyway, which is what the inherited value already says.
    """
    fallback = str(el.get("textAlign") or "start").lower()
    if not _FLEX_DISPLAY_RE.match(str(el.get("display") or "")):
        return fallback
    if "grid" in str(el.get("display") or ""):
        source = el.get("justifyItems") or "normal"
        if str(source).lower() in ("normal", "legacy", "stretch"):
            source = el.get("justifyContent") or "normal"
    elif "column" in str(el.get("flexDirection") or ""):
        # In a column the cross axis is the horizontal one.
        source = el.get("alignItems") or "normal"
    else:
        source = el.get("justifyContent") or "normal"
    return _ALIGN_KEYWORDS.get(str(source).lower(), fallback)


def _fill_representable(el: Dict[str, Any]) -> bool:
    """Can the element's own background be drawn as a native PowerPoint fill?"""
    if not _has_gradient(el):
        return True
    return _linear_gradient(el.get("bgImage")) is not None


def _border_sides(el: Dict[str, Any]) -> List[Tuple[float, str]]:
    """Per-side (width, color), in Top/Right/Bottom/Left order.

    Falls back to the scalar top-side pair when the array is absent, so a box
    styled only on, say, its left and bottom edges (an L-shaped corner mark) is
    still recognized as painted instead of being written off as a bare container.
    """
    widths = el.get("borderWidths")
    colors = el.get("borderColors")
    if isinstance(widths, list) and isinstance(colors, list) and len(widths) == len(colors) == 4:
        return list(zip((float(w or 0) for w in widths), colors))
    return [(float(el.get("borderWidth", 0) or 0), el.get("borderColor") or "")] * 4


def _has_visible_border(el: Dict[str, Any]) -> bool:
    return any(
        w > 0 and _alpha(c) > MIN_VISIBLE_ALPHA for w, c in _border_sides(el)
    )


def _is_panel(el: Dict[str, Any]) -> bool:
    """A styled box worth reconstructing (native shape or its own picture)."""
    if _alpha(el.get("bgColor")) > MIN_VISIBLE_ALPHA:
        return True
    if (el.get("bgImage") or "none") != "none":
        return True
    if _has_visible_border(el):
        return True
    if (el.get("boxShadow") or "none") != "none":
        return True
    return False


def _shape_representable(el: Dict[str, Any]) -> bool:
    """Can python-pptx draw this box faithfully as a native rectangle?

    Uniform border, uniform corner radius, no rotation, no gradient fill and no
    pseudo-element decoration. Anything else keeps more fidelity as its own
    picture — still an independent, movable object in the deck. Gradients are
    excluded even when PowerPoint could express them: a two-stop approximation
    of a multi-stop CSS ramp loses to a pixel-exact raster, and an element that
    reaches here has no children to swallow.
    """
    if not el.get("uniformBorder", True):
        return False
    if not el.get("uniformRadius", True):
        return False
    if abs(el.get("rotation", 0.0)) >= MIN_ROTATION_DEG:
        return False
    if _has_gradient(el):
        return False
    if el.get("hasPseudo"):
        return False
    return True


def classify_elements(elements: List[Dict[str, Any]]) -> Dict[int, str]:
    """Deterministic four-way role assignment.

    ``text``    editable text box (its own fill is emitted as a shape behind it)
    ``shape``   native (rounded) rectangle, children drawn on top
    ``picture`` rasterized on its own transparent canvas -> standalone picture
    ``image``   not reconstructed; whatever ink it has stays in the backdrop
    """
    roles: Dict[int, str] = {}
    for el in elements:
        eid = int(el["id"])
        big = el.get("areaFrac", 1.0) >= MAX_SHAPE_AREA_FRACTION
        if el.get("hasDirectText") and not el.get("hasTextAncestor"):
            roles[eid] = "text"
        elif el.get("isSvg") or el.get("isImg"):
            # Images and inline SVG become their own picture, never a shared backdrop.
            roles[eid] = "image" if big else "picture"
        elif el.get("hasTextAncestor"):
            roles[eid] = "image"          # part of an ancestor's text box
        elif _is_panel(el) and not big:
            if el.get("hasRenderableDescendant"):
                # A card stays native so its children draw on top — but only if we
                # can actually paint its own fill. Rasterizing is not an option
                # (it would swallow the children), and flattening an unexpressible
                # gradient to one color turns a faint radial glow into a solid
                # slab, so that ink is left in the backdrop instead.
                roles[eid] = "shape" if _fill_representable(el) else "image"
            else:
                roles[eid] = "shape" if _shape_representable(el) else "picture"
        elif el.get("hasPseudo") and not el.get("hasRenderableDescendant") and not big:
            roles[eid] = "picture"        # decoration drawn purely via ::before/::after
        else:
            roles[eid] = "image"
    return roles


def llm_classify(
    elements: List[Dict[str, Any]],
    full_png: Path,
    fallback: Dict[int, str],
) -> Dict[int, str]:
    """Single LLM(F2) call that may only override the role per element id.

    Coordinates/styles are never sourced from the model. Any missing or invalid
    id falls back to the deterministic decision. Returns ``fallback`` unchanged if
    credentials are absent or the call fails.
    """
    try:
        from ..prompts.loaders import load_prompt
        from ..providers.llm import get_llm

        system = load_prompt("tools/html_to_pptx.md")
        if not system.strip():
            return fallback
        table = [
            {
                "id": el["id"],
                "tag": el["tag"],
                "class": el["className"],
                "w": round(el["w"], 1),
                "h": round(el["h"], 1),
                "hasDirectText": el["hasDirectText"],
                "hasTextAncestor": el["hasTextAncestor"],
                "isSvg": el["isSvg"],
                "text": (el["innerText"] or "")[:60],
                "bgColor": el["bgColor"],
                "bgImage": "gradient" if _has_gradient(el) else "none",
                "borderRadius": round(el["borderRadius"], 1),
                "borderWidth": round(el["borderWidth"], 1),
                "default": fallback.get(int(el["id"])),
            }
            for el in elements
        ]
        user = (
            "Assign a PPTX role to every element id using the attached rendering "
            "as visual context. Roles: 'text' (editable text box), 'shape' (native "
            "rectangle), 'picture' (standalone rasterized image) or 'image' (leave "
            "in the backdrop). Reply with ONLY a JSON object mapping id (string) "
            "-> role.\n\nELEMENTS:\n" + json.dumps(table, ensure_ascii=False)
        )
        content = get_llm().vision(system, user, [str(full_png)])
        parsed = _extract_json_object(content)
        if not isinstance(parsed, dict):
            return fallback
        merged = dict(fallback)
        for key, role in parsed.items():
            try:
                eid = int(key)
            except (TypeError, ValueError):
                continue
            if role in {"text", "shape", "picture", "image"} and eid in merged:
                merged[eid] = role
        return merged
    except Exception:
        return fallback


def _extract_json_object(text: str) -> Any:
    if not text:
        return None
    stripped = text.strip()
    if stripped.startswith("```"):
        stripped = stripped.strip("`")
        stripped = re.sub(r"^json\s*", "", stripped, flags=re.IGNORECASE).strip()
    start = stripped.find("{")
    end = stripped.rfind("}")
    if start == -1 or end == -1 or end <= start:
        return None
    try:
        return json.loads(stripped[start : end + 1])
    except json.JSONDecodeError:
        return None


# Blend modes whose purpose is to make the white a raster sits on disappear.
MULTIPLY_BLEND_MODES = {"multiply", "darken"}


def bake_multiply_to_alpha(png_path: Path) -> bool:
    """Turn a ``mix-blend-mode: multiply`` raster into a plain transparent one.

    White-backed linework relies on multiply to make its background vanish
    against the paper. PowerPoint has no multiply for pictures, and an isolated
    shot blends against nothing, so the white would simply show.

    Multiply and alpha compositing agree when the ink is neutral: for backdrop
    ``d``, ``src*d`` equals ``a*C + (1-a)*d`` with ``C = 0`` and ``a = 1-src``.
    Taking ``m = min(r,g,b)`` as the amount of white to remove generalizes that
    to colored ink — exact over a white backdrop, and exact everywhere once the
    ink is gray, which is what a ``grayscale()`` filter leaves behind.

    Returns whether anything was written.
    """
    try:
        import numpy as np
        from PIL import Image
    except ImportError:
        return False

    with Image.open(png_path) as handle:
        image = handle.convert("RGBA")
    data = np.asarray(image).astype(np.float32) / 255.0
    rgb, alpha = data[:, :, :3], data[:, :, 3]

    white = rgb.min(axis=2)                       # the neutral part to knock out
    keep = 1.0 - white
    # Un-premultiply the ink against the white it sat on; where there is no ink
    # left the color is irrelevant because the alpha below is zero.
    safe = np.where(keep > 1e-4, keep, 1.0)[:, :, None]
    ink = np.clip((rgb - white[:, :, None]) / safe, 0.0, 1.0)

    out = np.empty_like(data)
    out[:, :, :3] = ink
    out[:, :, 3] = alpha * keep
    Image.fromarray((out * 255.0 + 0.5).astype(np.uint8), "RGBA").save(png_path)
    return True


def _resolve_canvas(page, canvas_selector: Optional[str]) -> str:
    """Pick the canvas selector: explicit wins, else first that matches exactly once."""
    if canvas_selector:
        count = page.locator(canvas_selector).count()
        if count != 1:
            raise RuntimeError(
                f"expected exactly one {canvas_selector}, found {count}"
            )
        return canvas_selector
    for selector in CANVAS_SELECTORS:
        if page.locator(selector).count() == 1:
            return selector
    raise RuntimeError(
        "no canvas found; pass canvas_selector explicitly "
        f"(tried {', '.join(CANVAS_SELECTORS)})"
    )


def capture_scene(
    index_path: Path,
    out_dir: Path,
    *,
    canvas_selector: Optional[str] = None,
    use_llm: bool = False,
    scale: int = 2,
    ready_selector: Optional[str] = None,
    prepare_js: Optional[str] = None,
    viewport: Optional[Tuple[int, int]] = None,
) -> Dict[str, Any]:
    """Render, extract, classify, rasterize every picture, and shoot the backdrop.

    ``ready_selector`` waits for a page that signals when it has settled (an
    intro animation, say) before anything is measured. ``prepare_js`` then runs
    in the page just before the canvas is resolved, so a caller can strip page
    chrome or defeat responsive scaling without this module knowing about any
    particular page. ``viewport`` pins the window size for pages that lay
    themselves out against it.
    """
    from playwright.sync_api import sync_playwright

    out_dir.mkdir(parents=True, exist_ok=True)
    parts_dir = out_dir / "parts"
    parts_dir.mkdir(parents=True, exist_ok=True)
    full_png = out_dir / "full.png"
    background_png = out_dir / "background.png"

    with _serve_preview(index_path) as url, sync_playwright() as playwright:
        browser_args = ["--no-sandbox"] if os.environ.get("HTML_TO_PPTX_NO_SANDBOX") == "1" else []
        browser = playwright.chromium.launch(args=browser_args)
        try:
            vw, vh = viewport or (1920, 1080)
            page = browser.new_page(
                viewport={"width": vw, "height": vh},
                device_scale_factor=scale,
            )
            page.goto(url, wait_until="load", timeout=30000)
            try:
                page.evaluate(
                    "async () => { if (document.fonts && document.fonts.ready)"
                    " { await document.fonts.ready; } }"
                )
            except Exception:
                pass
            if ready_selector:
                page.wait_for_selector(ready_selector, timeout=30000)
            page.wait_for_timeout(400)
            if prepare_js:
                page.evaluate(prepare_js)
                page.wait_for_timeout(200)

            selector = _resolve_canvas(page, canvas_selector)
            canvas = page.locator(selector)
            box = canvas.bounding_box() or {}
            cw = int(math.ceil(box.get("width") or 1920))
            ch = int(math.ceil(box.get("height") or 1080))
            if viewport is None:
                # Size the viewport to the canvas so nothing is clipped or scrolled.
                # Skipped when the caller pinned one: a page that scales itself to
                # fit the window would just shrink again in response.
                page.set_viewport_size({"width": cw, "height": ch})
                page.wait_for_timeout(150)

            extracted = page.evaluate(_EXTRACT_JS, selector)
            if not extracted:
                raise RuntimeError(f"{selector} not found during extraction")
            elements: List[Dict[str, Any]] = extracted["elements"]
            canvas_dims = extracted["canvas"]

            # Crisp full-canvas render used for color sampling.
            canvas.screenshot(path=str(full_png))

            roles = classify_elements(elements)
            if use_llm:
                roles = llm_classify(elements, full_png, roles)

            by_id = {int(el["id"]): el for el in elements}
            shape_ids = [i for i, r in roles.items() if r == "shape"]
            text_ids = [i for i, r in roles.items() if r == "text"]
            picture_ids = [i for i, r in roles.items() if r == "picture"]

            # Deterministic per-textbox shrink so text fits its box under font
            # substitution, independent of the opening app's autofit behavior.
            fit_items = [
                {"id": el["id"], "w": el["w"], "h": el["h"],
                 "family": el.get("fontFamily") or ""}
                for el in elements
                if int(el["id"]) in set(text_ids)
            ]
            fit_scales = page.evaluate(_MEASURE_FIT_JS, fit_items) if fit_items else {}
            for el in elements:
                el["fitScale"] = float(fit_scales.get(str(el["id"]), 1.0))

            # --- each picture on its own transparent canvas -------------------
            shot = page.viewport_size or {"width": cw, "height": ch}
            vw, vh = shot["width"], shot["height"]
            rasterized: Dict[int, str] = {}
            for eid in sorted(picture_ids):
                el = by_id[eid]
                # Shoot only the visible slice, so a picture hanging out of a crop
                # is cropped in the deck exactly as the page crops it.
                ox = _clamp(el.get("visVX", el["vx"]), 0.0, float(vw))
                oy = _clamp(el.get("visVY", el["vy"]), 0.0, float(vh))
                clip = {
                    "x": ox,
                    "y": oy,
                    "width": min(el.get("visW", el["w"]), vw - ox),
                    "height": min(el.get("visH", el["h"]), vh - oy),
                }
                # Chromium refuses sub-pixel clips; such a sliver is invisible anyway.
                if clip["width"] < 1.0 or clip["height"] < 1.0:
                    continue
                page.evaluate(_ISOLATE_JS, {"canvasSel": selector, "id": eid})
                part = parts_dir / f"el_{eid}.png"
                page.screenshot(path=str(part), clip=clip, omit_background=True)
                if str(el.get("mixBlendMode", "normal")).lower() in MULTIPLY_BLEND_MODES:
                    bake_multiply_to_alpha(part)
                rasterized[eid] = str(part)
            if picture_ids:
                page.evaluate(_RESET_ISOLATION_JS, selector)
                page.wait_for_timeout(80)

            # --- residual backdrop: only what we could not rebuild -------------
            page.evaluate(
                _NEUTRALIZE_JS,
                {
                    "shapeIds": shape_ids,
                    "textIds": text_ids,
                    "pictureIds": list(rasterized.keys()),
                },
            )
            page.wait_for_timeout(120)
            canvas.screenshot(path=str(background_png))
        finally:
            browser.close()

    for el in elements:
        el["partPng"] = rasterized.get(int(el["id"]), "")

    scene = {
        "canvas": canvas_dims,
        "scale": scale,
        "selector": selector,
        "background_png": str(background_png),
        "full_png": str(full_png),
        "elements": elements,
        "roles": {str(k): v for k, v in roles.items()},
    }
    (out_dir / "scene.json").write_text(
        json.dumps(scene, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return scene


# --------------------------------------------------------------------------- #
# PPTX builder
# --------------------------------------------------------------------------- #


def slide_size_emu(canvas_width: float, canvas_height: float) -> Tuple[int, int]:
    """Slide size preserving the canvas aspect ratio, long edge = LONG_SIDE_IN.

    1920x1080 -> the familiar 13.333x7.5in 16:9 deck; 1067x1600 -> portrait. The
    scale stays uniform on both axes, so artwork is never stretched.
    """
    cw = max(1.0, float(canvas_width))
    ch = max(1.0, float(canvas_height))
    if cw >= ch:
        win, hin = LONG_SIDE_IN, LONG_SIDE_IN * ch / cw
    else:
        hin, win = LONG_SIDE_IN, LONG_SIDE_IN * cw / ch
    return int(round(win * EMU_PER_INCH)), int(round(hin * EMU_PER_INCH))


def _sample_box_color(img: Any, el: Dict[str, Any], scale: int) -> Tuple[int, int, int]:
    """Median color of a ring just inside the element's edges.

    Sampling the exact center lands on a glyph whenever the box holds text,
    which used to paint whole pills the color of their own lettering.
    """
    if el.get("clipped"):
        bx, by, bw, bh = el["visX"], el["visY"], el["visW"], el["visH"]
    else:
        bx, by, bw, bh = el["x"], el["y"], el["w"], el["h"]
    x0 = int(bx * scale)
    y0 = int(by * scale)
    x1 = int((bx + bw) * scale)
    y1 = int((by + bh) * scale)
    x0, y0 = max(0, x0), max(0, y0)
    x1, y1 = min(img.width, x1), min(img.height, y1)
    if x1 - x0 < 3 or y1 - y0 < 3:
        px = max(0, min(img.width - 1, (x0 + x1) // 2))
        py = max(0, min(img.height - 1, (y0 + y1) // 2))
        pixel = img.getpixel((px, py))
        return (pixel, pixel, pixel) if isinstance(pixel, int) else tuple(pixel[:3])

    inset_x = max(1, (x1 - x0) // 10)
    inset_y = max(1, (y1 - y0) // 10)
    xs = range(x0 + inset_x, x1 - inset_x, max(1, (x1 - x0) // 24))
    ys = range(y0 + inset_y, y1 - inset_y, max(1, (y1 - y0) // 24))
    samples: List[Tuple[int, int, int]] = []
    for x in xs:
        for y in (y0 + inset_y, y1 - inset_y - 1):
            samples.append(_pixel_rgb(img, x, y))
    for y in ys:
        for x in (x0 + inset_x, x1 - inset_x - 1):
            samples.append(_pixel_rgb(img, x, y))
    if not samples:
        return _pixel_rgb(img, (x0 + x1) // 2, (y0 + y1) // 2)
    samples.sort(key=lambda c: c[0] + c[1] + c[2])
    return samples[len(samples) // 2]


def _pixel_rgb(img: Any, x: int, y: int) -> Tuple[int, int, int]:
    x = max(0, min(img.width - 1, x))
    y = max(0, min(img.height - 1, y))
    pixel = img.getpixel((x, y))
    if isinstance(pixel, int):
        return pixel, pixel, pixel
    return int(pixel[0]), int(pixel[1]), int(pixel[2])


def _clamp(value: float, low: float, high: float) -> float:
    return max(low, min(high, value))


def build_pptx(scene: Dict[str, Any], out_path: Path) -> Dict[str, Any]:
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Pt

    canvas = scene["canvas"]
    scale = int(scene.get("scale", 1))
    roles: Dict[int, str] = {int(k): v for k, v in scene["roles"].items()}
    by_id = {int(el["id"]): el for el in scene["elements"]}

    cw = float(canvas["width"])
    ch = float(canvas["height"])
    slide_w, slide_h = slide_size_emu(cw, ch)
    emu_x = slide_w / cw
    emu_y = slide_h / ch
    # Font/line sizes MUST scale by the same canvas->slide factor as positions.
    # The canvas (e.g. 1920px) is mapped onto 13.333in, so a CSS px is smaller
    # than 1/96in on the slide; using the naive 96dpi px->pt (0.75) makes every
    # font ~1.5x too big for its box. pt = px * EMU_per_px / EMU_per_pt.
    px_to_pt = emu_x / EMU_PER_PT

    def _ex(px: float) -> int:
        return int(round(_clamp(px, 0, cw) * emu_x))

    def _ey(px: float) -> int:
        return int(round(_clamp(px, 0, ch) * emu_y))

    def _ew(px: float) -> int:
        return max(1, int(round(px * emu_x)))

    def _eh(px: float) -> int:
        return max(1, int(round(px * emu_y)))

    align_map = {
        "left": PP_ALIGN.LEFT,
        "start": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "end": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    def _set_alpha(fill: Any, rgb: Tuple[int, int, int], alpha: float) -> None:
        fill.solid()
        fill.fore_color.rgb = RGBColor(*rgb)
        if alpha < 1.0:
            solid = fill._xPr.find(qn("a:solidFill"))
            srgb = solid.find(qn("a:srgbClr")) if solid is not None else None
            if srgb is not None:
                node = srgb.makeelement(
                    qn("a:alpha"),
                    {"val": str(int(_clamp(alpha, 0, 1) * 100000))},
                )
                srgb.append(node)

    def _set_gradient(
        fill: Any,
        angle: float,
        stops: List[Tuple[int, int, int, float]],
        opacity: float,
    ) -> None:
        """Two-stop linear gradient matching the CSS one end to end.

        PowerPoint's UI angle runs counter-clockwise from pointing right, while
        CSS measures clockwise from pointing up.
        """
        fill.gradient()
        slots = fill.gradient_stops
        for slot, rgba, position in (
            (slots[0], stops[0], 0.0),
            (slots[-1], stops[-1], 1.0),
        ):
            slot.color.rgb = RGBColor(*rgba[:3])
            slot.position = position
            alpha = rgba[3] * opacity
            if alpha < 1.0:
                srgb = slot._gs.find(qn("a:srgbClr"))
                if srgb is not None:
                    srgb.append(
                        srgb.makeelement(
                            qn("a:alpha"),
                            {"val": str(int(_clamp(alpha, 0, 1) * 100000))},
                        )
                    )
        fill.gradient_angle = (450.0 - angle) % 360.0

    def _set_run_color(run: Any, rgb: Tuple[int, int, int], alpha: float) -> None:
        """Run color with alpha. python-pptx exposes only opaque RGB, but a faded
        watermark word set fully opaque reads as a solid headline."""
        run.font.color.rgb = RGBColor(*rgb)
        if alpha >= 1.0:
            return
        rpr = run._r.get_or_add_rPr()
        solid = rpr.find(qn("a:solidFill"))
        srgb = solid.find(qn("a:srgbClr")) if solid is not None else None
        if srgb is not None:
            srgb.append(
                srgb.makeelement(
                    qn("a:alpha"), {"val": str(int(_clamp(alpha, 0, 1) * 100000))}
                )
            )

    def _set_font(run: Any, name: str) -> None:
        run.font.name = name
        rpr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            node = rpr.find(qn(tag))
            if node is None:
                node = rpr.makeelement(qn(tag), {})
                rpr.append(node)
            node.set("typeface", name)

    def _set_tracking(run: Any, points: float) -> None:
        """Letter spacing, in OOXML's hundredths of a point.

        Wide tracking is load-bearing in poster typography — a 29px word set at
        19px tracking collapses into a different design without it.
        """
        if abs(points) < 0.01:
            return
        run._r.get_or_add_rPr().set("spc", str(int(round(points * 100))))

    prs = Presentation()
    prs.slide_width = Emu(slide_w)
    prs.slide_height = Emu(slide_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Bottom layer: only what could not be reconstructed (usually page background).
    slide.shapes.add_picture(
        scene["background_png"], Emu(0), Emu(0), width=Emu(slide_w), height=Emu(slide_h)
    )

    full = Image.open(scene["full_png"]).convert("RGB")

    def _box(el: Dict[str, Any]) -> Tuple[float, float, float, float]:
        """Placement box: the visible slice once clipping ancestors are applied."""
        if el.get("clipped"):
            return el["visX"], el["visY"], el["visW"], el["visH"]
        return el["x"], el["y"], el["w"], el["h"]

    def _copacity(el: Dict[str, Any]) -> float:
        try:
            return _clamp(float(el.get("contentOpacity", 1.0) or 1.0), 0.0, 1.0)
        except (TypeError, ValueError):
            return 1.0

    def _cscale(el: Dict[str, Any]) -> float:
        """Scale that ancestor transforms already applied to the rendered box.

        Boxes are measured post-transform but font/border sizes are not, so both
        have to be multiplied through to stay in proportion.
        """
        try:
            value = float(el.get("contentScale", 1.0) or 1.0)
        except (TypeError, ValueError):
            return 1.0
        return value if value > 0 else 1.0

    def _add_box(el: Dict[str, Any]) -> Any:
        """Native (rounded) rectangle honoring the element's rotation."""
        bx, by, bw, bh = _box(el)
        rot = float(el.get("rotation", 0.0))
        if abs(rot) >= MIN_ROTATION_DEG:
            # The measured box is the axis-aligned bounds of a tilted rectangle,
            # so it is wider and taller than the rectangle itself. PowerPoint
            # rotates about the center: place the unrotated layout box there.
            cs_ = _cscale(el)
            bw, bh = (
                float(el.get("layoutW") or bw) * cs_,
                float(el.get("layoutH") or bh) * cs_,
            )
            bx, by = bx + (_box(el)[2] - bw) / 2.0, by + (_box(el)[3] - bh) / 2.0
        radius = el.get("borderRadius", 0) * _cscale(el)
        auto = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 1 else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(auto, _ex(bx), _ey(by), _ew(bw), _eh(bh))
        if radius > 1:
            short_side = max(1.0, min(bw, bh))
            shape.adjustments[0] = _clamp(radius / short_side, 0, 0.5)
        bg = _parse_rgba(el.get("bgColor"))
        border = _parse_rgba(el.get("borderColor"))
        outlined = _has_visible_border(el)
        op = _copacity(el)
        if bg and bg[3] > MIN_VISIBLE_ALPHA:
            _set_alpha(shape.fill, bg[:3], bg[3] * op)
        elif _has_gradient(el):
            gradient = _linear_gradient(el.get("bgImage"))
            if gradient:
                _set_gradient(shape.fill, gradient[0], gradient[1], op)
            else:
                _set_alpha(shape.fill, _sample_box_color(full, el, scale), op)
        elif outlined:
            # Outline-only box (hollow circle, hairline frame): filling it with a
            # sampled color turns a ring into a disc.
            shape.fill.background()
        else:
            _set_alpha(shape.fill, _sample_box_color(full, el, scale), 0.9 * op)
        if outlined and el.get("uniformBorder", True):
            shape.line.color.rgb = RGBColor(*border[:3])
            shape.line.width = Pt(max(0.5, el["borderWidth"] * _cscale(el) * px_to_pt))
        else:
            # Non-uniform borders (a colored top accent over gray sides) cannot be
            # one PPTX outline; each visible edge is drawn separately below.
            shape.line.fill.background()
        shape.shadow.inherit = False
        shape.text_frame.text = ""
        rot = float(el.get("rotation", 0.0))
        if abs(rot) >= MIN_ROTATION_DEG:
            shape.rotation = rot
        return shape

    def _add_border_sides(el: Dict[str, Any]) -> int:
        """Draw each visible edge of a non-uniform border as its own thin bar."""
        widths = el.get("borderWidths") or []
        colors = el.get("borderColors") or []
        if len(widths) != 4 or len(colors) != 4:
            return 0
        cs = _cscale(el)
        widths = [w * cs for w in widths]
        x, y, w, h = _box(el)
        geometry = (
            (x, y, w, widths[0]),                      # top
            (x + w - widths[1], y, widths[1], h),      # right
            (x, y + h - widths[2], w, widths[2]),      # bottom
            (x, y, widths[3], h),                      # left
        )
        made = 0
        for (bx, by, bw_, bh_), width, color in zip(geometry, widths, colors):
            rgba = _parse_rgba(color)
            if width <= 0 or not rgba or rgba[3] <= MIN_VISIBLE_ALPHA:
                continue
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, _ex(bx), _ey(by), _ew(bw_), _eh(bh_)
            )
            _set_alpha(bar.fill, rgba[:3], rgba[3] * _copacity(el))
            bar.line.fill.background()
            bar.shadow.inherit = False
            bar.text_frame.text = ""
            made += 1
        return made

    shape_count = text_count = picture_count = 0
    # Stacking order: the per-element paint path measured in the page, which
    # nests z-index inside its own stacking context. Falling back to a flat
    # z-index would let a positioned container's fill bury its own children.
    ordered = sorted(
        by_id.values(),
        key=lambda e: (
            list(e.get("paintPath") or [int(e.get("zIndex", 0) or 0)]),
            int(e["id"]),
        ),
    )

    for el in ordered:
        eid = int(el["id"])
        role = roles.get(eid)

        if role == "shape":
            _add_box(el)
            shape_count += 1
            if not el.get("uniformBorder", True):
                shape_count += _add_border_sides(el)

        elif role == "picture":
            part = el.get("partPng")
            if not part or not Path(part).exists():
                continue
            bx, by, bw, bh = _box(el)
            slide.shapes.add_picture(
                part, _ex(bx), _ey(by), width=_ew(bw), height=_eh(bh)
            )
            picture_count += 1

        elif role == "text":
            text = (el.get("innerText") or "").strip()
            if not text:
                continue
            # A badge/pill owns a fill *and* copy: emit the fill, then the words.
            if _is_panel(el):
                _add_box(el)
                shape_count += 1
                if not el.get("uniformBorder", True):
                    shape_count += _add_border_sides(el)

            cs = _cscale(el)
            tx, ty, tw, th = _box(el)
            rot = float(el.get("rotation", 0.0))
            if abs(rot) >= MIN_ROTATION_DEG:
                # PowerPoint rotates about the center, so place the *unrotated*
                # box centered on the measured (rotated) bounding box.
                lw = float(el.get("layoutW") or el["w"]) * cs
                lh = float(el.get("layoutH") or el["h"]) * cs
                cx = tx + tw / 2.0
                cy = ty + th / 2.0
                box = slide.shapes.add_textbox(
                    _ex(cx - lw / 2.0), _ey(cy - lh / 2.0), _ew(lw), _eh(lh)
                )
                box.rotation = rot
            else:
                box = slide.shapes.add_textbox(_ex(tx), _ey(ty), _ew(tw), _eh(th))

            tf = box.text_frame
            tf.word_wrap = True
            # CSS vertical writing-mode has a direct OOXML counterpart; without it
            # a vertical column of CJK glyphs reflows into horizontal lines.
            vertical = str(el.get("writingMode", "")).startswith("vertical")
            if vertical:
                tf._bodyPr.set("vert", "eaVert")
            # Shrink font to fit the original box so foreign-font substitution
            # cannot wrap text into extra lines that overlap neighbors.
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.margin_left = tf.margin_right = Emu(0)
            tf.margin_top = tf.margin_bottom = Emu(0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            base = _parse_rgba(el.get("color"))
            if not base or base[3] < 0.1:
                fallback_rgb = _gradient_first_color(el.get("bgImage")) or _sample_box_color(
                    full, el, scale
                )
                fallback_alpha = 1.0
            else:
                fallback_rgb = base[:3]
                fallback_alpha = base[3]

            # Only pay the anti-overflow margin when the text actually had to be
            # shrunk; when the real family is available the size stays exact.
            fit_scale = float(el.get("fitScale", 1.0))
            if fit_scale < 1.0:
                fit_scale *= 0.98
            el_size = el.get("fontSize", 16)
            align = align_map.get(_effective_align(el), PP_ALIGN.LEFT)
            # Everything measured in CSS px shares one conversion to points.
            type_scale = cs * px_to_pt * fit_scale
            el_family = str(el.get("fontFamily") or "")
            el_italic = str(el.get("fontStyle", "")) in ("italic", "oblique")
            el_tracking = float(el.get("letterSpacing", 0) or 0)
            line_height = float(el.get("lineHeight", 0) or 0)

            text_opacity = _copacity(el)

            def _run_color(value: Optional[str]) -> Tuple[Tuple[int, int, int], float]:
                parsed = _parse_rgba(value)
                if parsed and parsed[3] >= 0.1:
                    return parsed[:3], parsed[3] * text_opacity
                return fallback_rgb, fallback_alpha * text_opacity

            rich = el.get("richText") or []
            if not rich:
                rich = [[{"text": ln}] for ln in text.split("\n")]

            for i, line in enumerate(rich):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.alignment = align
                # In vertical writing-mode CSS line-height sizes the *column*, not
                # the advance between glyphs; mapping it to OOXML line spacing
                # pushes a single column of CJK into two.
                if line_height > 0 and not vertical:
                    para.line_spacing = Pt(line_height * type_scale)
                for seg in line:
                    seg_text = seg.get("text", "")
                    if seg_text == "":
                        continue
                    run = para.add_run()
                    run.text = seg_text
                    seg_size = float(seg.get("size") or el_size)
                    run.font.size = Pt(max(1.0, seg_size * type_scale))
                    try:
                        weight = int(float(seg.get("weight", el.get("fontWeight", 400))))
                    except (TypeError, ValueError):
                        weight = 700 if str(seg.get("weight")) == "bold" else 400
                    run.font.bold = weight >= BOLD_WEIGHT_THRESHOLD
                    run.font.italic = bool(seg.get("italic", el_italic))
                    _set_run_color(run, *_run_color(seg.get("color")))
                    tracking = seg.get("spacing")
                    _set_tracking(
                        run,
                        (el_tracking if tracking is None else float(tracking)) * type_scale,
                    )
                    family = str(seg.get("family") or el_family)
                    if not family:
                        # No concrete family survived the CSS chain; keep the old
                        # two-font guess so CJK does not land in a Latin face.
                        family = "Microsoft YaHei" if _CJK_RE.search(seg_text) else "Arial"
                    _set_font(run, family)
            text_count += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return {
        "shape_count": shape_count,
        "text_count": text_count,
        "picture_count": picture_count,
    }


def verify_pptx(out_path: Path, expected: Dict[str, Any]) -> Dict[str, Any]:
    """Deterministic sanity check (replaces a repair loop)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(out_path))
    slide = prs.slides[0]
    pictures = text_boxes = shapes = 0
    for shp in slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures += 1
        elif shp.has_text_frame and shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text_boxes += 1
        else:
            shapes += 1
    report = {
        "pictures": pictures,
        "text_boxes": text_boxes,
        "shapes": shapes,
        "expected_shapes": expected.get("shape_count", 0),
        "expected_text": expected.get("text_count", 0),
        "expected_pictures": expected.get("picture_count", 0) + 1,  # + backdrop
    }
    report["ok"] = (
        pictures >= expected.get("picture_count", 0) + 1
        and text_boxes >= expected.get("text_count", 0)
        and shapes >= expected.get("shape_count", 0)
    )
    return report


def build_editable_pptx(
    index_path: Path,
    out_dir: Path,
    *,
    out_name: str = "slide.pptx",
    canvas_selector: Optional[str] = None,
    use_llm: bool = False,
    scale: int = 2,
    ready_selector: Optional[str] = None,
    prepare_js: Optional[str] = None,
    viewport: Optional[Tuple[int, int]] = None,
) -> Dict[str, Any]:
    """Full pipeline entrypoint. Returns a manifest dict with the pptx path."""
    index_path = Path(index_path)
    out_dir = Path(out_dir)
    scene = capture_scene(
        index_path,
        out_dir,
        canvas_selector=canvas_selector,
        use_llm=use_llm,
        scale=scale,
        ready_selector=ready_selector,
        prepare_js=prepare_js,
        viewport=viewport,
    )
    pptx_path = out_dir / out_name
    counts = build_pptx(scene, pptx_path)
    verification = verify_pptx(pptx_path, counts)
    return {
        "pptx_path": str(pptx_path),
        "background_png": scene["background_png"],
        "scene_json": str(out_dir / "scene.json"),
        "canvas": scene["canvas"],
        "selector": scene["selector"],
        "element_count": len(scene["elements"]),
        **counts,
        "verification": verification,
    }
