#!/usr/bin/env python3
"""render_check —— 把生成的 PPTX 渲回图片，和源 HTML 比对并给出数值。

产出一张四联图：源 HTML ／ PPTX 渲回 ／ 每个独立对象的边界框 ／ 误差热力图。
第三联验证「元素是否真的各自独立」——框住的每一块在 PowerPoint 里都能单独选中。
第四联把误差最大的几块框出来，并在终端列出坐标与两侧色值，不用靠肉眼找问题在哪。

用法：
    python3 render_check.py <source.html> <out.pptx> [-o compare.png] [--selector .slide-canvas]
    python3 render_check.py <layers.html> <layers.pptx> --exploded

`--exploded` 用于 case 的图层展开图：截图前套用与转换器完全相同的页面预处理
（等动画就绪、删 overview、关响应式缩放），否则基准图里会多出总览图，比对没有意义。

依赖：LibreOffice（渲染 pptx）、PyMuPDF、Pillow、numpy、playwright(Python)、python-pptx。
"""
from __future__ import annotations

import argparse
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

CANVAS_SELECTORS = (".slide-canvas", ".poster-canvas", "[data-canvas-width]", "#poster")
SOFFICE_CANDIDATES = (
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "/usr/bin/soffice",
    "/usr/bin/libreoffice",
)
OUTLINE_COLORS = {"picture": (232, 62, 62), "text": (32, 120, 235), "shape": (28, 168, 92)}
# Mean channel distance above which a pixel counts as visibly wrong rather than
# as antialiasing or a font-hinting difference.
DIFF_THRESHOLD = 28
# Side of the coarse cell the diff is pooled into before regions are grown. Small
# enough to separate two neighbouring cards, large enough to bridge glyph gaps.
DIFF_CELL = 24


def _browser_args() -> list[str]:
    return ["--no-sandbox"] if os.environ.get("HTML_TO_PPTX_NO_SANDBOX") == "1" else []


def _find_soffice() -> str:
    for candidate in (os.environ.get("SOFFICE_PATH"), *SOFFICE_CANDIDATES):
        if candidate and Path(candidate).is_file():
            return candidate
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    sys.stderr.write(
        "找不到 LibreOffice（用于把 pptx 渲成图片）。\n"
        "  macOS: brew install --cask libreoffice\n"
        "  或用 SOFFICE_PATH 环境变量指定可执行文件路径。\n"
    )
    sys.exit(1)


def _load_exploded():
    """The展开图 entry point, for its page-prep contract."""
    import importlib.util

    path = Path(__file__).with_name("exploded_to_pptx.py")
    spec = importlib.util.spec_from_file_location("_exploded", path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _shoot_exploded(html: Path, out_png: Path) -> None:
    """Baseline shot of a layers board under the converter's own page prep.

    Shooting the page as it loads would capture the overview tile and a
    half-finished reveal animation, so the comparison would measure the
    difference between two different pictures.
    """
    from playwright.sync_api import sync_playwright

    exploded = _load_exploded()
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(args=_browser_args())
        page = browser.new_page(viewport={"width": 1600, "height": 1000})
        page.goto(html.resolve().as_uri(), wait_until="load", timeout=30000)
        try:
            page.wait_for_selector(exploded.READY, timeout=30000)
        except Exception:  # noqa: BLE001
            pass
        info = page.evaluate(exploded.PROBE_JS)
        page.set_viewport_size(
            {"width": max(1, int(info["width"])), "height": max(1, int(info["height"]))}
        )
        page.wait_for_timeout(400)
        page.evaluate(f"() => ({exploded.PREPARE_JS})({{'keepOverview': false}})")
        page.wait_for_timeout(300)
        page.locator(exploded.BOARD).screenshot(path=str(out_png))
        browser.close()


def _shoot_html(html: Path, selector: str | None, out_png: Path) -> None:
    from playwright.sync_api import sync_playwright

    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(args=_browser_args())
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.goto(html.resolve().as_uri(), wait_until="load", timeout=30000)
        try:
            page.evaluate(
                "async () => { if (document.fonts && document.fonts.ready)"
                " { await document.fonts.ready; } }"
            )
        except Exception:  # noqa: BLE001
            pass
        page.wait_for_timeout(400)
        picked = selector
        if not picked:
            for candidate in CANVAS_SELECTORS:
                if page.locator(candidate).count() == 1:
                    picked = candidate
                    break
        target = page.locator(picked) if picked else page
        box = target.bounding_box() if picked else None
        if box:
            page.set_viewport_size(
                {"width": max(1, int(box["width"])), "height": max(1, int(box["height"]))}
            )
            page.wait_for_timeout(150)
        target.screenshot(path=str(out_png))
        browser.close()


def _render_pptx(pptx: Path, work: Path) -> Path:
    soffice = _find_soffice()
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(work), str(pptx)],
        check=True,
        capture_output=True,
        timeout=300,
    )
    pdf = work / f"{pptx.stem}.pdf"
    if not pdf.is_file():
        sys.stderr.write(f"LibreOffice 未能生成 PDF: {pdf}\n")
        sys.exit(1)
    import fitz

    doc = fitz.open(pdf)
    png = work / "render.png"
    doc[0].get_pixmap(dpi=150).save(png)
    doc.close()
    return png


def _outline_objects(render_png: Path, pptx: Path, out_png: Path) -> int:
    from PIL import Image, ImageDraw
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    image = Image.open(render_png).convert("RGB")
    prs = Presentation(str(pptx))
    slide = prs.slides[0]
    sx = image.width / prs.slide_width
    sy = image.height / prs.slide_height
    draw = ImageDraw.Draw(image)
    count = 0
    for shape in slide.shapes:
        # The full-bleed fallback backdrop is not an independently editable object.
        if shape.width / prs.slide_width > 0.97 and shape.height / prs.slide_height > 0.97:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            kind = "picture"
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            kind = "text"
        else:
            kind = "shape"
        draw.rectangle(
            [
                shape.left * sx,
                shape.top * sy,
                (shape.left + shape.width) * sx,
                (shape.top + shape.height) * sy,
            ],
            outline=OUTLINE_COLORS[kind],
            width=3,
        )
        count += 1
    image.save(out_png)
    return count


def _worst_regions(bad, top: int):
    """Group flagged cells into regions, largest first.

    A plain breadth-first merge over the coarse grid; connected-component
    labelling would pull in SciPy for a few hundred cells.
    """
    rows, cols = bad.shape
    seen = [[False] * cols for _ in range(rows)]
    regions = []
    for r0 in range(rows):
        for c0 in range(cols):
            if not bad[r0][c0] or seen[r0][c0]:
                continue
            seen[r0][c0] = True
            queue = [(r0, c0)]
            cells = []
            while queue:
                r, c = queue.pop()
                cells.append((r, c))
                for dr, dc in ((1, 0), (-1, 0), (0, 1), (0, -1)):
                    nr, nc = r + dr, c + dc
                    if 0 <= nr < rows and 0 <= nc < cols and bad[nr][nc] and not seen[nr][nc]:
                        seen[nr][nc] = True
                        queue.append((nr, nc))
            rs = [c[0] for c in cells]
            cs = [c[1] for c in cells]
            regions.append((len(cells), min(cs), min(rs), max(cs) + 1, max(rs) + 1))
    regions.sort(reverse=True)
    return regions[:top]


def _diff_report(gt_png: Path, render_png: Path, out_png: Path, top: int = 5):
    """Pixel distance between the two panels, plus where the worst of it sits."""
    import numpy as np
    from PIL import Image, ImageDraw

    gt = Image.open(gt_png).convert("RGB")
    rendered = Image.open(render_png).convert("RGB").resize(gt.size, Image.LANCZOS)
    gt_a = np.asarray(gt, dtype=float)
    rd_a = np.asarray(rendered, dtype=float)
    delta = np.abs(gt_a - rd_a).mean(axis=2)

    heat = Image.fromarray(
        np.uint8(np.clip(delta * 3, 0, 255))
    ).convert("RGB")
    draw = ImageDraw.Draw(heat)

    rows = max(1, delta.shape[0] // DIFF_CELL)
    cols = max(1, delta.shape[1] // DIFF_CELL)
    pooled = delta[: rows * DIFF_CELL, : cols * DIFF_CELL]
    pooled = pooled.reshape(rows, DIFF_CELL, cols, DIFF_CELL).mean(axis=(1, 3))
    bad = pooled > DIFF_THRESHOLD

    found = []
    for _, c0, r0, c1, r1 in _worst_regions(bad, top):
        x0, y0, x1, y1 = c0 * DIFF_CELL, r0 * DIFF_CELL, c1 * DIFF_CELL, r1 * DIFF_CELL
        draw.rectangle([x0, y0, x1, y1], outline=(255, 90, 40), width=3)
        found.append(
            (
                (x0, y0, x1 - x0, y1 - y0),
                tuple(int(v) for v in gt_a[y0:y1, x0:x1].mean(axis=(0, 1))),
                tuple(int(v) for v in rd_a[y0:y1, x0:x1].mean(axis=(0, 1))),
            )
        )
    heat.save(out_png)
    return float(delta.mean()), float((delta > DIFF_THRESHOLD).mean() * 100), found


def _compose(panels: list[tuple[str, Path]], out_png: Path) -> None:
    from PIL import Image, ImageDraw

    height = 860
    images = []
    for label, path in panels:
        img = Image.open(path).convert("RGB")
        images.append((label, img.resize((round(img.width * height / img.height), height))))
    pad, bar = 14, 28
    total = sum(i.width for _, i in images) + pad * (len(images) + 1)
    canvas = Image.new("RGB", (total, height + bar + pad * 2), (28, 30, 34))
    draw = ImageDraw.Draw(canvas)
    x = pad
    for label, img in images:
        draw.text((x, 8), label, fill=(240, 240, 240))
        canvas.paste(img, (x, bar))
        x += img.width + pad
    canvas.save(out_png)


def main() -> None:
    ap = argparse.ArgumentParser(description="PPTX 渲回图片，与源 HTML 三联比对")
    ap.add_argument("html", help="源 slide/poster HTML")
    ap.add_argument("pptx", help="to_pptx.py 产出的 .pptx")
    ap.add_argument("-o", "--out", help="输出对比图（默认 pptx 同目录 <name>.compare.png）")
    ap.add_argument("--selector", default=None, help="画布选择器（默认自动识别）")
    ap.add_argument(
        "--exploded",
        action="store_true",
        help="源文件是 case 的图层展开图 layers.html，截图前套用转换器同款页面预处理",
    )
    args = ap.parse_args()

    html = Path(args.html).resolve()
    pptx = Path(args.pptx).resolve()
    for path in (html, pptx):
        if not path.is_file():
            sys.stderr.write(f"找不到文件: {path}\n")
            sys.exit(2)
    out = Path(args.out).resolve() if args.out else pptx.with_suffix(".compare.png")

    work = Path(tempfile.mkdtemp(prefix="render_check_"))
    gt = work / "gt.png"
    if args.exploded:
        _shoot_exploded(html, gt)
    else:
        _shoot_html(html, args.selector, gt)
    render = _render_pptx(pptx, work)
    outlined = work / "outlined.png"
    count = _outline_objects(render, pptx, outlined)
    heat = work / "heat.png"
    mean, over, regions = _diff_report(gt, render, heat)

    _compose(
        [
            ("源 HTML", gt),
            ("PPTX 渲回", render),
            (f"{count} 个独立可编辑对象（红=图片 蓝=文本框 绿=形状）", outlined),
            (f"误差热力图　平均 {mean:.2f}　超阈值 {over:.2f}%", heat),
        ],
        out,
    )
    print(f"✅ {out}")
    print(f"   独立对象 {count} 个（不含整页兜底背景）")
    print(f"   平均像素差 {mean:.2f}　超阈值占比 {over:.2f}%")
    if regions:
        print("   误差最大的区域（源图坐标 x,y,w,h）：")
        for (x, y, w, h), gt_rgb, rd_rgb in regions:
            print(f"     ({x},{y},{w},{h})  源 rgb{gt_rgb}  →  pptx rgb{rd_rgb}")
    else:
        print("   没有超出阈值的区域")


if __name__ == "__main__":
    main()
